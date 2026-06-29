import os
import subprocess
import tempfile
import zipfile
from io import BytesIO
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pypdf import PdfWriter, PdfReader
from PIL import Image
import fitz
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
import pdfplumber
import openpyxl
import platform
import shutil

app = Flask(__name__)
CORS(app)

# Use correct LibreOffice path depending on the OS
if platform.system() == 'Windows':
    LIBREOFFICE_PATH = r"C:\Program Files\LibreOffice\program\soffice.exe"
else:
    # Use absolute path for Linux so os.path.exists() works
    LIBREOFFICE_PATH = "/usr/bin/libreoffice"

@app.route('/api/merge', methods=['POST'])
def merge_pdfs():
    if 'files' not in request.files:
        return jsonify({"error": "No files part"}), 400
    
    files = request.files.getlist('files')
    if not files or files[0].filename == '':
        return jsonify({"error": "No selected files"}), 400

    merger = PdfWriter()
    try:
        for file in files:
            if file.filename.endswith('.pdf'):
                merger.append(file)
        
        output_io = BytesIO()
        merger.write(output_io)
        merger.close()
        output_io.seek(0)
        
        return send_file(output_io, mimetype='application/pdf', as_attachment=True, download_name='merged.pdf')
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/convert/pdf-to-image', methods=['POST'])
def convert_pdf_to_image():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '' or not file.filename.endswith('.pdf'):
        return jsonify({"error": "Invalid or no PDF file selected"}), 400

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
        file.save(temp_pdf.name)
        temp_pdf_path = temp_pdf.name

    try:
        # Open PDF with PyMuPDF
        doc = fitz.open(temp_pdf_path)
        
        memory_file = BytesIO()
        with zipfile.ZipFile(memory_file, 'w') as zf:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=150)
                img_data = pix.tobytes("png")
                zf.writestr(f"page_{i+1}.png", img_data)
        
        doc.close()
        memory_file.seek(0)
        
        baseName = os.path.splitext(file.filename)[0]
        return send_file(memory_file, mimetype='application/zip', as_attachment=True, download_name=f"{baseName}.zip")
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        if os.path.exists(temp_pdf_path):
            os.remove(temp_pdf_path)

@app.route('/api/security/protect', methods=['POST'])
def protect_pdf():
    if 'file' not in request.files or 'password' not in request.form:
        return jsonify({"error": "Missing file or password"}), 400
    
    file = request.files['file']
    password = request.form['password']
    
    if file.filename == '' or not file.filename.endswith('.pdf'):
        return jsonify({"error": "Invalid or no PDF file selected"}), 400

    try:
        reader = PdfReader(file)
        writer = PdfWriter()
        
        for page in reader.pages:
            writer.add_page(page)
            
        writer.encrypt(password)
        
        output_io = BytesIO()
        writer.write(output_io)
        output_io.seek(0)
        
        return send_file(output_io, mimetype='application/pdf', as_attachment=True, download_name='protected.pdf')
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/convert/office-to-pdf', methods=['POST'])
def convert_office_to_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No file selected"}), 400
        
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in ['.docx', '.pptx', '.xlsx']:
        return jsonify({"error": "Unsupported file format. Use .docx, .pptx, or .xlsx"}), 400

    if not os.path.exists(LIBREOFFICE_PATH):
        return jsonify({"error": "LibreOffice not found on the server at the configured path."}), 500

    with tempfile.TemporaryDirectory() as temp_dir:
        input_path = os.path.join(temp_dir, file.filename)
        file.save(input_path)
        
        try:
            # Run LibreOffice headless conversion
            process = subprocess.run(
                [LIBREOFFICE_PATH, "--headless", "--convert-to", "pdf", input_path, "--outdir", temp_dir],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE
            )
            
            if process.returncode != 0:
                return jsonify({"error": f"Conversion failed: {process.stderr.decode()}"}), 500
                
            pdf_filename = os.path.splitext(file.filename)[0] + ".pdf"
            output_pdf_path = os.path.join(temp_dir, pdf_filename)
            
            if not os.path.exists(output_pdf_path):
                return jsonify({"error": "Conversion failed: output file not generated"}), 500
                
            # Read the generated PDF into memory so we can return it before temp_dir is deleted
            with open(output_pdf_path, 'rb') as f:
                pdf_data = f.read()
                
            output_io = BytesIO(pdf_data)
            return send_file(output_io, mimetype='application/pdf', as_attachment=True, download_name=pdf_filename)
            
        except Exception as e:
            return jsonify({"error": str(e)}), 500

@app.route('/api/convert/image-to-pdf', methods=['POST'])
def convert_image_to_pdf():
    if 'files' not in request.files:
        return jsonify({"error": "No files part"}), 400
    
    files = request.files.getlist('files')
    if not files or files[0].filename == '':
        return jsonify({"error": "No selected files"}), 400

    image_list = []
    try:
        for file in files:
            ext = os.path.splitext(file.filename)[1].lower()
            if ext not in ['.png', '.jpg', '.jpeg']:
                continue
            
            img = Image.open(file.stream)
            # PDF conversion requires RGB mode (convert transparent backgrounds)
            if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
                # Paint white background behind transparency
                background = Image.new("RGB", img.size, (255, 255, 255))
                background.paste(img, mask=img.split()[3] if img.mode == 'RGBA' else None)
                img = background
            else:
                img = img.convert('RGB')
            image_list.append(img)
        
        if not image_list:
            return jsonify({"error": "No valid PNG, JPG, or JPEG images were provided."}), 400
        
        output_io = BytesIO()
        first_img = image_list[0]
        first_img.save(output_io, format='PDF', save_all=True, append_images=image_list[1:])
        output_io.seek(0)
        
        return send_file(output_io, mimetype='application/pdf', as_attachment=True, download_name='converted_images.pdf')
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/split', methods=['POST'])
def split_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '' or not file.filename.endswith('.pdf'):
        return jsonify({"error": "Invalid or no PDF file selected"}), 400

    split_type = request.form.get('split_type', 'split_all')
    page_range = request.form.get('page_range', '')

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
        file.save(temp_pdf.name)
        temp_pdf_path = temp_pdf.name

    try:
        reader = PdfReader(temp_pdf_path)
        total_pages = len(reader.pages)
        if total_pages == 0:
            return jsonify({"error": "The PDF file is empty."}), 400

        base_name = os.path.splitext(file.filename)[0]

        if split_type == 'extract_range':
            if not page_range.strip():
                return jsonify({"error": "Page range must be specified for extraction."}), 400
            
            pages_to_extract = []
            try:
                parts = page_range.split(',')
                for part in parts:
                    part = part.strip()
                    if '-' in part:
                        start, end = part.split('-')
                        pages_to_extract.extend(range(int(start) - 1, int(end)))
                    else:
                        pages_to_extract.append(int(part) - 1)
            except Exception:
                return jsonify({"error": "Invalid page range format. Use formats like '1-3, 5'."}), 400

            valid_pages = [p for p in pages_to_extract if 0 <= p < total_pages]
            if not valid_pages:
                return jsonify({"error": "No valid pages found in the specified range."}), 400

            writer = PdfWriter()
            for p in valid_pages:
                writer.add_page(reader.pages[p])

            output_io = BytesIO()
            writer.write(output_io)
            writer.close()
            output_io.seek(0)
            return send_file(output_io, mimetype='application/pdf', as_attachment=True, download_name=f"{base_name}_split.pdf")

        else:
            memory_file = BytesIO()
            with zipfile.ZipFile(memory_file, 'w') as zf:
                for i in range(total_pages):
                    writer = PdfWriter()
                    writer.add_page(reader.pages[i])
                    page_io = BytesIO()
                    writer.write(page_io)
                    writer.close()
                    zf.writestr(f"page_{i+1}.pdf", page_io.getvalue())
            
            memory_file.seek(0)
            return send_file(memory_file, mimetype='application/zip', as_attachment=True, download_name=f"{base_name}_split.zip")

    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        if os.path.exists(temp_pdf_path):
            os.remove(temp_pdf_path)

@app.route('/api/convert/pdf-to-docx', methods=['POST'])
def convert_pdf_to_docx():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    file = request.files['file']
    if file.filename == '' or not file.filename.endswith('.pdf'):
        return jsonify({"error": "Invalid or no PDF file selected"}), 400

    with tempfile.TemporaryDirectory() as temp_dir:
        pdf_path = os.path.join(temp_dir, "input.pdf")
        docx_path = os.path.join(temp_dir, "output.docx")
        file.save(pdf_path)
        try:
            cv = Converter(pdf_path)
            cv.convert(docx_path, start=0, end=None)
            cv.close()
            
            with open(docx_path, 'rb') as f:
                docx_data = f.read()
            output_io = BytesIO(docx_data)
            base_name = os.path.splitext(file.filename)[0]
            return send_file(output_io, mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document', as_attachment=True, download_name=f"{base_name}.docx")
        except Exception as e:
            return jsonify({"error": str(e)}), 500

@app.route('/api/convert/pdf-to-pptx', methods=['POST'])
def convert_pdf_to_pptx():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    file = request.files['file']
    if file.filename == '' or not file.filename.endswith('.pdf'):
        return jsonify({"error": "Invalid or no PDF file selected"}), 400

    with tempfile.TemporaryDirectory() as temp_dir:
        pdf_path = os.path.join(temp_dir, "input.pdf")
        file.save(pdf_path)
        try:
            prs = Presentation()
            doc = fitz.open(pdf_path)
            if len(doc) == 0:
                return jsonify({"error": "PDF is empty"}), 400

            first_page = doc[0]
            rect = first_page.rect
            prs.slide_width = Inches(rect.width / 72.0)
            prs.slide_height = Inches(rect.height / 72.0)

            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=150)
                img_path = os.path.join(temp_dir, f"page_{i}.png")
                pix.save(img_path)

                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

            doc.close()
            pptx_path = os.path.join(temp_dir, "output.pptx")
            prs.save(pptx_path)

            with open(pptx_path, 'rb') as f:
                pptx_data = f.read()
            output_io = BytesIO(pptx_data)
            base_name = os.path.splitext(file.filename)[0]
            return send_file(output_io, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation', as_attachment=True, download_name=f"{base_name}.pptx")
        except Exception as e:
            return jsonify({"error": str(e)}), 500

@app.route('/api/convert/pdf-to-xlsx', methods=['POST'])
def convert_pdf_to_xlsx():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    file = request.files['file']
    if file.filename == '' or not file.filename.endswith('.pdf'):
        return jsonify({"error": "Invalid or no PDF file selected"}), 400

    with tempfile.TemporaryDirectory() as temp_dir:
        pdf_path = os.path.join(temp_dir, "input.pdf")
        file.save(pdf_path)
        try:
            wb = openpyxl.Workbook()
            wb.remove(wb.active)

            with pdfplumber.open(pdf_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    ws = wb.create_sheet(title=f"Page {i+1}")
                    tables = page.extract_tables()

                    if tables:
                        current_row = 1
                        for table in tables:
                            for row in table:
                                for col_idx, cell in enumerate(row):
                                    ws.cell(row=current_row, column=col_idx+1, value=cell)
                                current_row += 1
                            current_row += 2
                    else:
                        text = page.extract_text()
                        if text:
                            for row_idx, line in enumerate(text.split('\n')):
                                ws.cell(row=row_idx+1, column=1, value=line)
                        else:
                            ws.cell(row=1, column=1, value="[No extractable text or tables on this page]")

            xlsx_path = os.path.join(temp_dir, "output.xlsx")
            wb.save(xlsx_path)

            with open(xlsx_path, 'rb') as f:
                xlsx_data = f.read()
            output_io = BytesIO(xlsx_data)
            base_name = os.path.splitext(file.filename)[0]
            return send_file(output_io, mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', as_attachment=True, download_name=f"{base_name}.xlsx")
        except Exception as e:
            return jsonify({"error": str(e)}), 500

@app.route('/api/compress', methods=['POST'])
def compress_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '' or not file.filename.endswith('.pdf'):
        return jsonify({"error": "Invalid or no PDF file selected"}), 400

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
        file.save(temp_pdf.name)
        temp_pdf_path = temp_pdf.name

    try:
        doc = fitz.open(temp_pdf_path)
        output_io = BytesIO()
        # Save compressed to memory
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as out_pdf:
            out_pdf_path = out_pdf.name
        
        doc.save(out_pdf_path, deflate=True, garbage=4)
        doc.close()
        
        with open(out_pdf_path, 'rb') as f:
            pdf_data = f.read()
        output_io = BytesIO(pdf_data)
        
        if os.path.exists(out_pdf_path):
            os.remove(out_pdf_path)
            
        base_name = os.path.splitext(file.filename)[0]
        return send_file(output_io, mimetype='application/pdf', as_attachment=True, download_name=f"{base_name}_compressed.pdf")
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        if os.path.exists(temp_pdf_path):
            os.remove(temp_pdf_path)

@app.route('/api/edit', methods=['POST'])
def edit_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '' or not file.filename.endswith('.pdf'):
        return jsonify({"error": "Invalid or no PDF file selected"}), 400

    watermark_text = request.form.get('watermark', 'DocFlow Edited')

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
        file.save(temp_pdf.name)
        temp_pdf_path = temp_pdf.name

    try:
        doc = fitz.open(temp_pdf_path)
        for page in doc:
            rect = page.rect
            # Add watermark text diagonally across the page
            page.insert_text(
                (rect.width / 4, rect.height / 2), 
                watermark_text, 
                fontsize=40, 
                color=(0.8, 0.8, 0.8), # Light grey
                overlay=True
            )
            
        output_io = BytesIO()
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as out_pdf:
            out_pdf_path = out_pdf.name
        
        doc.save(out_pdf_path)
        doc.close()
        
        with open(out_pdf_path, 'rb') as f:
            pdf_data = f.read()
        output_io = BytesIO(pdf_data)
        
        if os.path.exists(out_pdf_path):
            os.remove(out_pdf_path)
            
        base_name = os.path.splitext(file.filename)[0]
        return send_file(output_io, mimetype='application/pdf', as_attachment=True, download_name=f"{base_name}_edited.pdf")
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        if os.path.exists(temp_pdf_path):
            os.remove(temp_pdf_path)

if __name__ == '__main__':
    app.run(debug=True, port=5000)
